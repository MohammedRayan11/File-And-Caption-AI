import os
from PIL import Image
import fitz  # PyMuPDF for PDF text extraction
from pdf2docx import Converter
from moviepy.editor import VideoFileClip
from pydub import AudioSegment
import shutil
import logging
from docx import Document
from pptx import Presentation
import subprocess
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from docx2pdf import convert  # Correct import

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def ensure_output_directory():
    """Ensure the 'converted' directory exists."""
    if not os.path.exists('converted'):
        os.makedirs('converted')

def pptx_to_pdf(pptx_path):
    """Convert PowerPoint (PPTX) to PDF."""
    try:
        output_path = pptx_path.replace(".pptx", ".pdf")
        subprocess.run(["unoconv", "-f", "pdf", "-o", output_path, pptx_path], check=True)
        return output_path
    except Exception as e:
        logging.error(f"Error converting PPTX to PDF: {e}")
        raise

def cdr_to_png(cdr_path):
    """Convert CorelDRAW (CDR) to PNG using Inkscape."""
    try:
        output_path = cdr_path.replace(".cdr", ".png")
        subprocess.run(["inkscape", cdr_path, "--export-png", output_path], check=True)
        return output_path
    except Exception as e:
        logging.error(f"Error converting CDR to PNG: {e}")
        raise

def pdf_to_text(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        text = "\n".join(page.get_text() for page in doc)
        
        # Ensure the output directory exists
        os.makedirs('converted', exist_ok=True)
        
        # Construct the output path
        output_path = os.path.join('converted', os.path.basename(pdf_path).replace(".pdf", ".txt"))
        
        # Save the text to the output file
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        
        logging.info(f"Text file saved to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error converting PDF to text: {e}")
        raise

def docx_to_html(docx_path):
    """Convert DOCX to HTML."""
    try:
        doc = Document(docx_path)
        html_content = "<html><body>"

        for para in doc.paragraphs:
            html_content += f"<p>{para.text}</p>"

        html_content += "</body></html>"
        html_path = docx_path.replace(".docx", ".html")

        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html_content)

        return html_path
    except Exception as e:
        logging.error(f"Error converting DOCX to HTML: {e}")
        raise

def docx_to_text(docx_path):
    """Convert DOCX to text."""
    try:
        doc = Document(docx_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        text_path = docx_path.replace(".docx", ".txt")

        with open(text_path, "w", encoding="utf-8") as f:
            f.write(text)

        return text_path
    except Exception as e:
        logging.error(f"Error converting DOCX to text: {e}")
        raise

import win32com.client

import os
import shutil
import logging

def png_to_svg(png_path):
    try:
        # Ensure the input file exists
        if not os.path.exists(png_path):
            raise FileNotFoundError(f"Input file not found: {png_path}")

        # Define the output path
        output_path = png_path.replace(".png", ".svg")
        output_path = os.path.join('converted', os.path.basename(output_path))

        # Ensure the output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        # Execute the Inkscape command to trace the PNG and export as SVG
        command = ["inkscape", png_path, "--export-plain-svg", output_path]
        logging.info(f"Executing command: {command}")
        subprocess.run(command, check=True)

        # Verify that the file was saved
        if not os.path.exists(output_path):
            raise FileNotFoundError(f"Output file not found: {output_path}")

        logging.info(f"Successfully converted {png_path} to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error converting PNG to SVG: {e}")
        raise

def svg_to_cdr(svg_path):
    try:
        # Ensure the input file exists
        if not os.path.exists(svg_path):
            raise FileNotFoundError(f"Input file not found: {svg_path}")

        # Define the output path
        output_path = svg_path.replace(".svg", ".cdr")
        output_path = os.path.join('converted', os.path.basename(output_path))

        # Ensure the output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        # Copy the SVG file and rename it to CDR
        shutil.copy(svg_path, output_path)
        logging.info(f"Renamed {svg_path} to {output_path}")

        return output_path
    except Exception as e:
        logging.error(f"Error renaming SVG to CDR: {e}")
        raise
def png_to_cdr(png_path):
    """Convert PNG to CDR by first converting to SVG, then to CDR."""
    try:
        # Convert PNG to SVG
        svg_path = png_to_svg(png_path)
        
        # Convert SVG to CDR
        return svg_to_cdr(svg_path)
    except Exception as e:
        logging.error(f"Error converting PNG to CDR: {e}")
        raise
def jpg_to_cdr(jpg_path):
    """Convert JPG to CDR by first converting to PNG, then to SVG, then to CDR."""
    try:
        # Convert JPG to PNG
        png_path = jpg_path.replace(".jpg", ".png")
        with Image.open(jpg_path) as img:
            img.save(png_path)

        # Convert PNG to CDR
        return png_to_cdr(png_path)
    except Exception as e:
        logging.error(f"Error converting JPG to CDR: {e}")
        raise
def pdf_to_svg(pdf_path):
    """Convert PDF to SVG using pdf2svg."""
    try:
        # Ensure the output directory exists
        ensure_output_directory()

        # Define the output SVG path
        svg_path = pdf_path.replace(".pdf", ".svg")
        svg_path = os.path.join('converted', os.path.basename(svg_path))

        # Convert PDF to SVG using pdf2svg
        subprocess.run(["pdf2svg", pdf_path, svg_path], check=True)

        logging.info(f"Successfully converted {pdf_path} to {svg_path}")
        return svg_path
    except subprocess.CalledProcessError as e:
        logging.error(f"PDF to SVG conversion failed: {e}")
        raise
    except FileNotFoundError:
        logging.error("pdf2svg is not installed. Please install it to proceed.")
        raise
    except Exception as e:
        logging.error(f"Error converting PDF to SVG: {e}")
        raise
from docx2pdf import convert
import logging
import os

import pypandoc

import pythoncom  # Import pythoncom for COM initialization

import pypandoc

import pythoncom  # Import pythoncom for COM initialization

def docx_to_pdf(docx_path):
    try:
        # Ensure the output directory exists
        os.makedirs('converted', exist_ok=True)
        
        # Define the output PDF file path
        output_path = os.path.join('converted', os.path.basename(docx_path).replace(".docx", ".pdf"))
        
        # Convert DOCX to PDF using pypandoc
        pypandoc.convert_file(docx_path, 'pdf', outputfile=output_path)
        
        logging.info(f"PDF file saved to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error converting DOCX to PDF: {e}")
        raise
def txt_to_pdf(txt_path):
    try:
        output_path = txt_path.replace(".txt", ".pdf")
        output_path = os.path.join('converted', os.path.basename(output_path))
        
        c = canvas.Canvas(output_path, pagesize=letter)
        y = 750  # Start position for text

        with open(txt_path, "r", encoding="utf-8") as file:
            for line in file:
                c.drawString(50, y, line.strip())
                y -= 15  # Move down for the next line

        c.save()
        logging.info(f"PDF saved to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error converting TXT to PDF: {e}")
        raise

def pdf_to_cdr(pdf_path):
    """Convert PDF to CDR by first converting to SVG, then to CDR."""
    try:
        # Convert PDF to SVG
        svg_path = pdf_path.replace(".pdf", ".svg")
        subprocess.run(["pdf2svg", pdf_path, svg_path], check=True)
        
        # Convert SVG to CDR
        return svg_to_cdr(svg_path)
    except Exception as e:
        logging.error(f"Error converting PDF to CDR: {e}")
        raise
def cdr_to_png(cdr_path):
    """Convert CDR to PNG using Inkscape."""
    try:
        output_path = cdr_path.replace(".cdr", ".png")
        subprocess.run(["inkscape", cdr_path, "--export-png", output_path], check=True)
        return output_path
    except Exception as e:
        logging.error(f"Error converting CDR to PNG: {e}")
        raise

def cdr_to_pdf(cdr_path):
    """Convert CDR to PDF using Inkscape."""
    try:
        output_path = cdr_path.replace(".cdr", ".pdf")
        subprocess.run(["inkscape", cdr_path, "--export-pdf", output_path], check=True)
        return output_path
    except Exception as e:
        logging.error(f"Error converting CDR to PDF: {e}")
        raise

def cdr_to_svg(cdr_path):
    """Convert CDR to SVG using Inkscape."""
    try:
        output_path = cdr_path.replace(".cdr", ".svg")
        subprocess.run(["inkscape", cdr_path, "--export-plain-svg", output_path], check=True)
        return output_path
    except Exception as e:
        logging.error(f"Error converting CDR to SVG: {e}")
        raise
def html_to_pdf(html_path):
    """Convert HTML to PDF."""
    try:
        output_path = html_path.replace(".html", ".pdf")
        # Use an external tool like wkhtmltopdf for better HTML to PDF conversion
        subprocess.run(["wkhtmltopdf", html_path, output_path], check=True)
        return output_path
    except Exception as e:
        logging.error(f"Error converting HTML to PDF: {e}")
        raise

from pdfminer.high_level import extract_text_to_fp
from io import StringIO

def pdf_to_html(pdf_path):
    try:
        # Ensure the output directory exists
        os.makedirs('converted', exist_ok=True)
        
        # Define the output HTML file path
        output_path = os.path.join('converted', os.path.basename(pdf_path).replace(".pdf", ".html"))
        
        # Extract text from PDF and save as HTML
        with open(output_path, "wb") as output_file:
            extract_text_to_fp(open(pdf_path, "rb"), output_file, output_type="html")
        
        logging.info(f"HTML file saved to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error converting PDF to HTML: {e}")
        raise

import fitz  # PyMuPDF

def pdf_to_md(pdf_path):
    try:
        # Ensure the output directory exists
        os.makedirs('converted', exist_ok=True)
        
        # Define the output Markdown file path
        output_path = os.path.join('converted', os.path.basename(pdf_path).replace(".pdf", ".md"))
        
        # Extract text from PDF
        doc = fitz.open(pdf_path)
        markdown_content = ""
        
        for page in doc:
            text = page.get_text()
            # Format text as Markdown
            markdown_content += text + "\n\n"  # Add extra newlines for Markdown paragraphs
        
        # Write the Markdown content to the output file
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(markdown_content)
        
        logging.info(f"Markdown file saved to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error converting PDF to Markdown: {e}")
        raise
import fitz  # PyMuPDF
import pypandoc

from ebooklib import epub

def pdf_to_epub(pdf_path):
    try:
        # Ensure the output directory exists
        os.makedirs('converted', exist_ok=True)
        
        # Define the output EPUB file path
        output_path = os.path.join('converted', os.path.basename(pdf_path).replace(".pdf", ".epub"))
        
        # Extract text from PDF
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text() + "\n\n"

        # Create an EPUB book
        book = epub.EpubBook()

        # Set metadata
        book.set_title("Converted PDF")
        book.set_language("en")

        # Create a chapter
        chapter = epub.EpubHtml(title="Chapter 1", file_name="chap_01.xhtml", lang="en")
        chapter.content = f"<h1>Chapter 1</h1><p>{text}</p>"

        # Add chapter to the book
        book.add_item(chapter)

        # Define Table of Contents
        book.toc = (epub.Link("chap_01.xhtml", "Chapter 1", "chap_01"),)

        # Add default NCX and Nav files
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())

        # Define CSS style (optional)
        style = "BODY { font-family: Times, serif; }"
        nav_css = epub.EpubItem(uid="style_nav", file_name="style/nav.css", media_type="text/css", content=style)
        book.add_item(nav_css)

        # Create spine
        book.spine = ["nav", chapter]

        # Save the EPUB file
        epub.write_epub(output_path, book, {})
        
        logging.info(f"EPUB file saved to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error converting PDF to EPUB: {e}")
        raise
import fitz  # PyMuPDF
import pypandoc

import subprocess

def pdf_to_odt(pdf_path):
    try:
        # Ensure the output directory exists
        os.makedirs('converted', exist_ok=True)
        
        # Define the output ODT file path
        output_path = os.path.join('converted', os.path.basename(pdf_path).replace(".pdf", ".odt"))
        
        # Use LibreOffice to convert PDF to ODT
        subprocess.run(["soffice", "--convert-to", "odt", "--outdir", "converted", pdf_path], check=True)
        
        logging.info(f"ODT file saved to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error converting PDF to ODT: {e}")
        raise
from pdf2docx import Converter
import pypandoc

def pdf_to_rtf(pdf_path):
    try:
        # Ensure the output directory exists
        os.makedirs('converted', exist_ok=True)
        
        # Define the output RTF file path
        output_path = os.path.join('converted', os.path.basename(pdf_path).replace(".pdf", ".rtf"))
        
        # Step 1: Convert PDF to DOCX
        docx_path = os.path.join('converted', os.path.basename(pdf_path).replace(".pdf", ".docx"))
        cv = Converter(pdf_path)
        cv.convert(docx_path, start=0, end=None)
        cv.close()
        
        # Step 2: Convert DOCX to RTF
        pypandoc.convert_file(docx_path, 'rtf', outputfile=output_path)
        
        # Clean up the intermediate DOCX file
        os.remove(docx_path)
        
        logging.info(f"RTF file saved to {output_path}")
        return output_path
    except Exception as e:
        logging.error(f"Error converting PDF to RTF: {e}")
        raise
import subprocess
import os

import subprocess
import os
import pypandoc

def docx_to_md(docx_path):
    """Convert DOCX to Markdown (MD) format."""
    try:
        output_path = docx_path.replace(".docx", ".md")
        pypandoc.convert_file(docx_path, 'md', outputfile=output_path)
        return output_path
    except Exception as e:
        logging.error(f"Error converting DOCX to Markdown: {e}")
        raise
def pdf_to_cdr(pdf_path):
    """Convert PDF to CDR by first converting to SVG, then to CDR."""
    try:
        # Step 1: Convert PDF to SVG using pdftocairo
        svg_path = pdf_path.replace(".pdf", ".svg")
        subprocess.run(["pdftocairo", "-svg", pdf_path, svg_path], check=True)

        # Step 2: Convert SVG to CDR using Inkscape
        output_path = pdf_path.replace(".pdf", ".cdr")
        subprocess.run(["inkscape", svg_path, "--export-plain-svg", output_path], check=True)

        # Clean up the intermediate SVG file
        os.remove(svg_path)

        return output_path
    except subprocess.CalledProcessError as e:
        logging.error(f"PDF to CDR conversion failed: {e}")
        raise
    except FileNotFoundError:
        logging.error("pdftocairo or Inkscape is not installed. Please install them to proceed.")
        raise
    except Exception as e:
        logging.error(f"Error converting PDF to CDR: {e}")
        raise
import pypandoc

def docx_to_rtf(docx_path):
    """Convert DOCX to RTF format."""
    try:
        output_path = docx_path.replace(".docx", ".rtf")
        pypandoc.convert_file(docx_path, 'rtf', outputfile=output_path)
        return output_path
    except Exception as e:
        logging.error(f"Error converting DOCX to RTF: {e}")
        raise   
def docx_to_epub(docx_path):
    """Convert DOCX to EPUB format."""
    try:
        output_path = docx_path.replace(".docx", ".epub")
        pypandoc.convert_file(docx_path, 'epub', outputfile=output_path)
        return output_path
    except Exception as e:
        logging.error(f"Error converting DOCX to EPUB: {e}")
        raise   
def docx_to_odt(docx_path):
    """Convert DOCX to ODT format."""
    try:
        output_path = docx_path.replace(".docx", ".odt")
        pypandoc.convert_file(docx_path, 'odt', outputfile=output_path)
        return output_path
    except Exception as e:
        logging.error(f"Error converting DOCX to ODT: {e}")
        raise
def convert_file(input_path, output_format):
    """
    Converts a file from one format to another.
    :param input_path: Path to the input file.
    :param output_format: Desired output format.
    :return: Path to the converted file.
    """
    try:
        ensure_output_directory()
        
        filename = os.path.basename(input_path)
        output_filename = f"{os.path.splitext(filename)[0]}.{output_format}"
        output_path = os.path.join('converted', output_filename)
        logging.info(f"Attempting to convert {input_path} to {output_format}")
        logging.info(f"Output path: {output_path}")
        if input_path.lower().endswith(".pdf"):
            if output_format == "html":
                return pdf_to_html(input_path)
            elif output_format == "epub":
                return pdf_to_epub(input_path)
            elif output_format == "odt":
                return pdf_to_odt(input_path)
            elif output_format == "rtf":
                return pdf_to_rtf(input_path)
            elif output_format == "cdr":
                return pdf_to_cdr(input_path)
            elif output_format == "md":  # Add Markdown conversion
                return pdf_to_md(input_path)
            elif output_format == "docx":
                cv = Converter(input_path)
                cv.convert(output_path, start=0, end=None)  # Convert entire document
                cv.close()
                return output_path
            elif output_format == "txt":
                text_file = pdf_to_text(input_path)
                return text_file  # Return extracted text file
            else:
                raise ValueError(f"Unsupported PDF conversion format: {output_format}")
        # 📌 DOCX to PDF Conversion
        elif input_path.lower().endswith(".docx") and output_format == "pdf":
            try:
        # Initialize the COM library
                pythoncom.CoInitialize()

                convert(input_path, output_path)
                return output_path
            except Exception as e:
                logging.error(f"Error converting DOCX to PDF: {e}")
                raise
            finally:
        # Uninitialize the COM library
                pythoncom.CoUninitialize()
        
        elif output_format == "cdr":
            if input_path.lower().endswith(".png"):
                return png_to_cdr(input_path)
            elif input_path.lower().endswith((".jpg", ".jpeg")):
                return jpg_to_cdr(input_path)
            elif input_path.lower().endswith(".svg"):
                return svg_to_cdr(input_path)
            elif input_path.lower().endswith(".pdf"):
                return pdf_to_cdr(input_path)
            else:
                raise ValueError(f"Unsupported input format for CDR conversion: {input_path}")
        # 📌 PDF Conversions (FROM PDF)
        elif input_path.lower().endswith(".pdf"):
            if output_format == "docx":
                cv = Converter(input_path)
                cv.convert(output_path, start=0, end=None)  # Convert entire document
                cv.close()
                return output_path
            elif output_format == "txt":
                text_file = pdf_to_text(input_path)
                return text_file  # Return extracted text file
            else:
                raise ValueError(f"Unsupported PDF conversion format: {output_format}")
        elif input_path.lower().endswith(".svg") and output_format == "cdr":
            return svg_to_cdr(input_path)
        
            
        
        # 📌 Convert from CDR
        elif input_path.lower().endswith(".cdr"):
            if output_format == "png":
                return cdr_to_png(input_path)
            elif output_format == "pdf":
                return cdr_to_pdf(input_path)
            elif output_format == "svg":
                return cdr_to_svg(input_path)
            else:
                raise ValueError(f"Unsupported CDR conversion format: {output_format}")
        # 📌 TXT Conversions (FROM TXT)
        elif input_path.lower().endswith(".txt"):
            if output_format == "pdf":
                return txt_to_pdf(input_path)
            elif output_format == "docx":
                doc = Document()
                with open(input_path, "r", encoding="utf-8") as file:
                    for line in file:
                        doc.add_paragraph(line.strip())
                doc.save(output_path)
                logging.info(f"File saved to {output_path}")
                return output_path
            else:
                raise ValueError(f"Unsupported TXT conversion format: {output_format}")

        # 📌 HTML Conversions (FROM HTML)
        elif input_path.lower().endswith(".html"):
            if output_format == "pdf":
                return html_to_pdf(input_path)
            else:
                raise ValueError(f"Unsupported HTML conversion format: {output_format}")

        # 📌 Convert to CDR
        elif output_format == "cdr":
            if input_path.lower().endswith(".svg"):
                # Convert SVG to CDR
                return svg_to_cdr(input_path)
            elif input_path.lower().endswith((".png", ".jpg", ".jpeg")):
                # Convert PNG/JPG to SVG first, then to CDR
                svg_path = png_to_svg(input_path)
                return svg_to_cdr(svg_path)
            elif input_path.lower().endswith(".pdf"):
                # Convert PDF to SVG first, then to CDR
                svg_path = pdf_to_svg(input_path)  # Implement this function if needed
                return svg_to_cdr(svg_path)
            else:
                raise ValueError(f"Unsupported input format for CDR conversion: {input_path}")

        # 📌 PowerPoint (PPTX) Conversions
        elif input_path.lower().endswith(".pptx"):
            if output_format == "pdf":
                return pptx_to_pdf(input_path)
            elif output_format == "docx":
                # Convert PPTX to DOCX (requires custom logic or external tool)
                raise ValueError("PPTX to DOCX conversion is not yet supported.")
            elif output_format == "txt":
                # Extract text from PPTX
                prs = Presentation(input_path)
                text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(text)
                return output_path
            else:
                raise ValueError(f"Unsupported PPTX conversion format: {output_format}")

        # 📌 CorelDRAW (CDR) Conversions
        elif input_path.lower().endswith(".cdr"):
            if output_format in ["png", "jpg", "pdf", "svg"]:
                # Convert CDR to PNG first
                png_path = cdr_to_png(input_path)
                
                # If the output format is PNG, return the PNG file
                if output_format == "png":
                    return png_path
                
                # Convert PNG to other formats (JPG, PDF, SVG)
                img = Image.open(png_path)
                if output_format == "jpg":
                    img.convert("RGB").save(output_path, "JPEG")
                elif output_format == "pdf":
                    img.convert("RGB").save(output_path, "PDF")
                elif output_format == "svg":
                    # Convert PNG to SVG using Inkscape
                    subprocess.run(["inkscape", png_path, "--export-plain-svg", output_path], check=True)
                return output_path
            else:
                raise ValueError(f"Unsupported CDR conversion format: {output_format}")

        # 📌 Adobe Illustrator (AI) and Photoshop (PSD) Conversions
        elif input_path.lower().endswith((".ai", ".psd")):
            if output_format in ["png", "jpg", "pdf", "svg"]:
                img = Image.open(input_path)
                img.save(output_path)
                return output_path
            else:
                raise ValueError(f"Unsupported AI/PSD conversion format: {output_format}")

        # 📌 SVG Conversions
        elif input_path.lower().endswith(".svg"):
            if output_format in ["png", "jpg", "pdf", "ai"]:
                img = Image.open(input_path)
                img.save(output_path)
                return output_path
            else:
                raise ValueError(f"Unsupported SVG conversion format: {output_format}")

        # 📌 Document Conversions (TXT to DOCX, etc.)
        elif input_path.lower().endswith((".txt", ".docx")) and output_format in ["pdf", "docx", "txt"]:
            if output_format == "pdf" and input_path.endswith(".txt"):
                from reportlab.lib.pagesizes import letter
                from reportlab.pdfgen import canvas

                with open(input_path, "r", encoding="utf-8") as file:
                    content = file.readlines()

                c = canvas.Canvas(output_path, pagesize=letter)
                y = 750  # Start position for text

                for line in content:
                    c.drawString(50, y, line.strip())  # Write each line in the PDF
                    y -= 15  # Move down

                c.save()
                return output_path

            elif output_format == "docx" and input_path.endswith(".txt"):
                doc = Document()
                with open(input_path, "r", encoding="utf-8") as file:
                    for line in file:
                        doc.add_paragraph(line.strip())

                doc.save(output_path)
                return output_path

            elif output_format == "txt" and input_path.endswith(".docx"):
                return docx_to_text(input_path)
        elif input_path.lower().endswith(".docx") and output_format == "html":
            return docx_to_html(input_path)
        
        elif input_path.lower().endswith(".docx"):
            if output_format == "pdf":
                return docx_to_pdf(input_path)
            elif output_format == "txt":
                return docx_to_text(input_path)
            elif output_format == "html":
                return docx_to_html(input_path)
            elif output_format == "md":  # Add Markdown conversion
                return docx_to_md(input_path)
            elif output_format == "rtf":  # Add RTF conversion
                return docx_to_rtf(input_path)
            elif output_format == "odt":  # Add ODT conversion
                return docx_to_odt(input_path)
            elif output_format == "epub":  # Add EPUB conversion
                return docx_to_epub(input_path)
            else:
                raise ValueError(f"Unsupported DOCX conversion format: {output_format}")

        # 📌 Image Conversions
        elif output_format in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp', 'tiff', 'pdf']:
            img = Image.open(input_path)
            if output_format == "pdf":
                img.convert("RGB").save(output_path, "PDF")
            else:
                img.save(output_path)
            return output_path

        # 📌 Audio Conversions
        elif output_format in ['mp3', 'wav', 'ogg', 'flac', 'aac']:
            audio = AudioSegment.from_file(input_path)
            audio.export(output_path, format=output_format)
            return output_path

        # 📌 Video Conversions
        elif output_format in ['mp4', 'avi', 'mkv', 'mov', 'wmv', 'flv']:
            video = VideoFileClip(input_path)
            video.write_videofile(output_path, codec='libx264', audio_codec="aac")
            return output_path

        # 📌 Archive Conversions
        elif output_format in ['zip', 'tar', 'gz']:
            shutil.make_archive(output_path.replace(f".{output_format}", ""), output_format, os.path.dirname(input_path), filename)
            return output_path
        
    except Exception as e:
        logging.error(f"Error converting file: {e}")
        raise